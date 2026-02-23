"""
==============================================================================
CampusMind: RAG 기반 지능형 학습 보조 시스템
==============================================================================
이 앱은 사용자가 업로드한 강의 자료(PDF, PPT, Word, 이미지, 음성, 영상)와
연습 문제를 벡터 데이터베이스(FAISS)에 저장한 뒤,
GPT-4o를 통해 다음 기능을 제공합니다:
  1. 핵심 정리 (요약)
  2. 시각화 (마인드맵 / 스파이더 다이어그램)
  3. 플래시카드 (Q&A 암기 카드)
  4. 퀴즈 (4지선다 객관식)
  5. 오디오 브리핑 (TTS 낭독)
  6. AI 도우미 (RAG 기반 질의응답)

아키텍처: RAG (Retrieval-Augmented Generation)
  - 문서 → 청크 분할 → 임베딩 → FAISS 벡터 DB 저장
  - 질의 → 벡터 검색 → 관련 청크 검색 → GPT-4o 답변 생성
==============================================================================
"""

# ==============================================================================
# [1] 표준 라이브러리 임포트
# ==============================================================================
import streamlit as st   # Streamlit: Python으로 웹앱을 만들 수 있는 프레임워크
import os                # 운영체제 파일 경로, 파일 존재 여부 등 시스템 작업용
import tempfile          # 업로드된 파일을 디스크에 임시 저장하기 위한 모듈
import base64            # 이미지 파일을 OpenAI Vision API에 전송할 때 Base64 인코딩 사용
import re                # 정규표현식: JSON 파싱 전처리, 마크다운 제거 등에 사용
import json              # LLM이 반환한 JSON 문자열을 Python 딕셔너리로 파싱
import time              # (예약) 재시도 딜레이 등에 사용할 수 있는 시간 모듈

# ==============================================================================
# [2] AI / LangChain 라이브러리 임포트
# ==============================================================================

# LangChain 기반 OpenAI 연동 모듈
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
# - OpenAIEmbeddings: 텍스트를 벡터(숫자 배열)로 변환하는 임베딩 모델 (문서 저장/검색에 사용)
# - ChatOpenAI: GPT-4o 채팅 모델에 대한 LangChain 래퍼 (invoke()로 호출)

from langchain_community.vectorstores import FAISS
# FAISS (Facebook AI Similarity Search): 임베딩 벡터를 저장하고
# 코사인 유사도 기반으로 유사한 문서를 빠르게 검색해주는 벡터 데이터베이스

from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
# - PyPDFLoader: PDF 파일을 페이지 단위로 텍스트를 추출하는 로더
# - Docx2txtLoader: Word 문서(.docx/.doc)의 텍스트를 추출하는 로더

from langchain_core.documents import Document
# LangChain에서 문서 단위를 표현하는 기본 클래스
# page_content(텍스트)와 metadata(출처 정보 등)를 갖는 구조체

from langchain_text_splitters import RecursiveCharacterTextSplitter
# 긴 문서를 일정 크기(chunk_size)의 청크(조각)로 분할하는 텍스트 분할기
# 단락 → 문장 → 단어 순으로 재귀적으로 분할하여 문맥을 최대한 보존

from langchain.chains import RetrievalQA
# 검색 기반 QA 체인: 사용자 질문 → 벡터 DB 검색 → 검색된 문서를 컨텍스트로 하여 LLM 답변 생성

from langchain.prompts import PromptTemplate
# LLM에 전달할 프롬프트 템플릿을 구조화하는 클래스
# {context}, {question} 같은 변수를 포함한 문자열 템플릿을 생성

from langchain_community.callbacks import get_openai_callback
# LangChain 체인 실행 중 OpenAI API 토큰 사용량을 실시간으로 추적하는 콜백
# "with get_openai_callback() as cb:" 블록 안에서 cb.prompt_tokens, cb.completion_tokens 확인 가능

from pptx import Presentation
# python-pptx: PowerPoint 파일(.pptx)에서 슬라이드의 텍스트를 추출하기 위한 라이브러리

from openai import OpenAI
# OpenAI 공식 Python SDK: Whisper(음성→텍스트), TTS(텍스트→음성), Vision(이미지 분석) API 직접 호출에 사용

import graphviz  # Graphviz DOT 언어로 작성된 다이어그램을 렌더링하는 라이브러리 (Streamlit과 연동)

# ==============================================================================
# [3] 멀티미디어 라이브러리 임포트 (선택적)
# ==============================================================================
# MoviePy는 pip install moviepy 버전에 따라 임포트 경로가 다를 수 있어
# 최신 버전(moviepy 2.x)과 구 버전(moviepy 1.x)을 모두 시도하는 방어적 임포트 사용
try:
    from moviepy import VideoFileClip      # moviepy 2.x 이상 (최신 버전)
except ImportError:
    try:
        from moviepy.editor import VideoFileClip  # moviepy 1.x (구 버전)
    except ImportError:
        VideoFileClip = None  # 두 버전 모두 없으면 None으로 설정 → 영상 처리 기능 비활성화


# ==============================================================================
# [4] Streamlit 페이지 기본 설정
# ==============================================================================
# 브라우저 탭 제목, 레이아웃, 파비콘(아이콘) 설정
# layout="wide": 화면 전체 너비를 활용하는 넓은 레이아웃 (기본값은 "centered")
# ※ 이 설정은 반드시 st.write() 등 다른 Streamlit 함수 호출보다 먼저 실행되어야 함
st.set_page_config(page_title="CampusMind", layout="wide", page_icon="🧠")

# ==============================================================================
# [5] 커스텀 CSS 스타일 적용
# ==============================================================================
# unsafe_allow_html=True: Streamlit이 기본적으로 막아두는 HTML/CSS 직접 삽입을 허용
# 전체 앱의 폰트, 버튼 스타일, 탭 활성화 색상, 그래프 정렬 등을 커스터마이징
st.markdown("""
<style>
    /* Google Fonts에서 한국어 지원 폰트(Noto Sans KR)를 웹으로 불러옴 */
    @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@400;500;700&display=swap');
    
    /* 모든 텍스트 요소에 Noto Sans KR 폰트를 강제 적용 (!important로 기본값 덮어쓰기) */
    h1, h2, h3, h4, h5, h6, p, li, label, textarea, input, div { 
        font-family: 'Noto Sans KR', sans-serif !important; 
    }
    /* Material Icons 아이콘 폰트는 예외 처리 (아이콘이 깨지지 않도록) */
    .material-icons, .material-symbols-rounded {
        font-family: 'Material Icons' !important;
    }
    /* Streamlit 버튼: 폰트 적용 + 너비 100%로 확장 (사이드바 버튼이 꽉 차도록) */
    .stButton > button {
        font-family: 'Noto Sans KR', sans-serif !important;
        width: 100%;
    }
    /* 마크다운 컴포넌트에도 폰트 적용 */
    .stMarkdown {
        font-family: 'Noto Sans KR', sans-serif !important;
    }
    /* 탭 버튼에도 폰트 적용 */
    .stTabs button {
        font-family: 'Noto Sans KR', sans-serif !important;
    }
    /* 확장 패널(expander)에 테두리와 배경색 부여 → 시각적 구분 */
    .stExpander {
        border: 1px solid #e0e0e0;
        border-radius: 8px;
        background-color: #ffffff;
        margin-bottom: 10px;
    }
    /* 선택된 탭에 하늘색 배경 + 파란 상단 보더 + 굵은 글씨 적용 */
    .stTabs [data-baseweb="tab-list"] button[aria-selected="true"] {
        background-color: #e3f2fd;
        border-top: 3px solid #1976d2;
        color: #0d47a1;
        font-weight: bold;
    }
    /* Graphviz SVG 차트를 화면 너비에 맞게 꽉 채우고 가운데 정렬 */
    [data-testid="stGraphvizChart"] svg {
        max-width: 100% !important;
        height: auto !important;
        display: block;
        margin: 0 auto;
    }
</style>
""", unsafe_allow_html=True)


# ==============================================================================
# [6] UI 언어 팩 (한국어 / 영어 다국어 지원)
# ==============================================================================
# 앱에서 사용되는 모든 UI 문자열을 딕셔너리로 정의합니다.
# 이 구조 덕분에 "lang_opt" 변수 하나만 바꾸면 전체 UI 언어가 전환됩니다.
# 새로운 언어를 추가하려면 동일한 키(key) 구조로 새 딕셔너리를 추가하면 됩니다.
UI = {
    "Korean": {
        "title":          "🧠 CampusMind: 지능형 학습 보조 시스템",
        "credit":         "By 이충환",
        "caption":        "Architecture: RAG-based LLM Workflow",
        "sidebar_title":  "⚙️ 데이터 소스",
        "file_label_lec": "📚 강의 자료 (PDF, PPT, Word, 이미지, 음성, 영상 등)",
        "file_label_prob":"📝 연습 문제 (PDF, Word 등)",
        "apikey":         "OpenAI API 키",
        "btn_start":      "🚀 분석 시작",
        # 6개 탭의 레이블 (순서 중요: t1~t6에 대응)
        "tabs":           ["📝 핵심 정리", "🎨 시각화", "🃏 플래시카드", "🧩 퀴즈", "🎧 오디오 브리핑", "💬 AI 도우미"],
        "input_topic":    "주제 필터 (전부, 전체, 빈칸 시 전체 범위)",
        "ph_topic":       "예: '신경망' (전부, 전체, 빈칸 시 전체 범위)",
        "msg_proc":       "📥 데이터 처리 중...",
        "msg_ingest":     "읽는 중: ",
        "msg_done":       "✅ 분석 완료!",
        "msg_err_file":   "파일 처리 오류: ",
        "msg_nodata":     "데이터 없음.",
        "btn_gen":        "생성하기",
        "viz_types":      ["Mindmap", "Spider Diagram"],  # 시각화 유형 선택지
        "quiz_check":     "정답 확인",
        "quiz_correct":   "정답입니다! ⭕",
        "quiz_wrong":     "오답입니다. ❌",
        "quiz_exp":       "해설 보기",
        "target_lang":    "Korean",   # LLM 출력 언어 지시에 사용되는 실제 언어명
        "lbl_card_front": "질문",
        "lbl_card_back":  "정답",
        "audio_btn":      "🎙️ 오디오 브리핑 생성",
        "audio_warn":     "먼저 요약을 생성해주세요.",
        "spin_gen":       "생성 중...",
        "spin_viz":       "구조화 중...",
        "spin_audio":     "오디오 합성 중...",
        "err_viz":        "렌더링 오류. Graphviz가 설치되어 있는지 확인하세요.",
        "err_viz_debug":  "DOT 코드 확인 (디버깅)",
        "chat_ph":        "질문을 입력하세요...",
        # h_bullet, h_table, h_term: 요약 섹션 제목 (오디오 추출용 마커로도 사용됨)
        "h_bullet":       "1. 핵심 내용 요약",
        "h_table":        "2. 상세 요약 표",
        "h_term":         "3. 용어 정리",
        # h_th: 표의 열(column) 헤더 이름 리스트 [구분, 상세설명, 용어, 정의, 문맥]
        "h_th":           ["구분", "상세 설명", "용어", "정의", "문맥"],
        "err_json":       "데이터 생성 오류. 다시 시도해주세요."
    },
    "English": {
        "title":          "🧠 CampusMind: Intelligent Tutor System",
        "credit":         "By Choonghwan Lee",
        "caption":        "Architecture: RAG-based LLM Workflow",
        "sidebar_title":  "⚙️ Data Sources",
        "file_label_lec": "📚 Lecture Materials (PDF, PPT, Word, Image, Audio, Video)",
        "file_label_prob":"📝 Practice Problems (PDF, Word)",
        "apikey":         "OpenAI API Key",
        "btn_start":      "🚀 Analyze",
        "tabs":           ["📝 Summary", "🎨 Visuals", "🃏 Flashcards", "🧩 Quiz", "🎧 Audio Brief", "💬 AI Tutor"],
        "input_topic":    "Topic Filter (All, Everything, Blank for All Sections)",
        "ph_topic":       "e.g., 'Neural Networks' (All, Everything, Blank for All Sections)",
        "msg_proc":       "📥 Processing Data...",
        "msg_ingest":     "Ingesting: ",
        "msg_done":       "✅ Ready!",
        "msg_err_file":   "File Error: ",
        "msg_nodata":     "No data.",
        "btn_gen":        "Generate",
        "viz_types":      ["Mindmap", "Spider Diagram"],
        "quiz_check":     "Check Answer",
        "quiz_correct":   "Correct! ⭕",
        "quiz_wrong":     "Incorrect. ❌",
        "quiz_exp":       "Explanation",
        "target_lang":    "English",
        "lbl_card_front": "Question",
        "lbl_card_back":  "Answer",
        "audio_btn":      "🎙️ Generate Audio",
        "audio_warn":     "Generate summary first.",
        "spin_gen":       "Generating...",
        "spin_viz":       "Generating diagram...",
        "spin_audio":     "Synthesizing...",
        "err_viz":        "Rendering Error. Please check Graphviz installation.",
        "err_viz_debug":  "View DOT Code",
        "chat_ph":        "Ask a question...",
        "h_bullet":       "1. Key Highlights",
        "h_table":        "2. Detailed Summary Table",
        "h_term":         "3. Terminology",
        "h_th":           ["Category", "Detailed Content", "Term", "Definition", "Context"],
        "err_json":       "Generation Error. Please try again."
    }
}


# ==============================================================================
# [7] 공통 유틸리티: 전체 모드 판별
# ==============================================================================

# 사용자가 주제 필터에 입력할 수 있는 "전체 범위"를 의미하는 키워드 집합
# set 자료형을 사용하여 O(1) 시간복잡도로 포함 여부를 확인
# 이 집합을 수정하면 gen_summary, gen_diagram, gen_flashcards, gen_quiz 모두에 자동 반영됨 (Fix #7)
ALL_MODE_KEYWORDS = {"all", "전부", "전체", "everything", "total"}

def check_is_all_mode(topic: str) -> bool:
    """
    주어진 topic 문자열이 "전체 범위" 모드인지 판별합니다.

    동작 방식:
    - topic이 None이거나 공백이면 → True (전체 범위로 처리)
    - topic을 소문자로 변환한 뒤 ALL_MODE_KEYWORDS에 포함되면 → True
    - 그 외에는 → False (특정 주제 검색 모드)

    이 함수를 공통 헬퍼로 빼둔 이유:
    기존에는 4개 함수(gen_summary, gen_diagram, gen_flashcards, gen_quiz) 각각에
    동일한 if/else 블록이 중복되어 있었습니다. 키워드 목록 변경 시 4군데를 수정해야 했지만,
    이 함수 하나만 수정하면 됩니다. (Fix #7)
    """
    return not topic or topic.strip().lower() in ALL_MODE_KEYWORDS


# ==============================================================================
# [8] API 비용 추적 시스템
# ==============================================================================
# OpenAI API는 사용량(토큰 수, 글자 수, 재생 시간)에 따라 과금됩니다.
# 이 섹션에서는 각 API 호출마다 비용을 계산하고,
# Streamlit 세션에 누적하여 사이드바에 실시간으로 표시합니다.

# --- 2026년 1월 기준 OpenAI 공식 요금표 (USD) ---
# 출처: https://platform.openai.com/pricing
# ※ 요금은 변경될 수 있으므로 실제 청구 전에 공식 페이지를 반드시 확인하세요.
PRICING = {
    # GPT-4o: 텍스트 생성 모델 (요약, 퀴즈, 플래시카드, 시각화, AI 튜터에 사용)
    # 입력 토큰(프롬프트)과 출력 토큰(응답)을 별도로 과금
    "gpt-4o": {
        "input":  2.50 / 1_000_000,   # $2.50 / 1M input tokens  → 토큰당 $0.0000025
        "output": 10.00 / 1_000_000,  # $10.00 / 1M output tokens → 토큰당 $0.00001
    },
    # TTS-1: 텍스트를 음성으로 변환하는 표준 모델 (오디오 브리핑 기능에 사용)
    # 글자(character) 수 기준으로 과금 (공백, 구두점 포함)
    "tts-1": {
        "chars":  15.00 / 1_000_000,  # $15.00 / 1M characters → 글자당 $0.000015
    },
    # Whisper-1: 음성을 텍스트로 변환하는 전사(transcription) 모델
    # 오디오 재생 시간(분) 기준으로 과금 (초 단위는 올림 처리)
    "whisper-1": {
        "per_min": 0.006,             # $0.006 / minute → 분당 0.6센트
    },
    # text-embedding-3-small: 텍스트를 벡터로 변환하는 임베딩 모델
    # 2026년 기준 권장 임베딩 모델 (구 ada-002 대비 성능 향상 + 가격 80% 절감)
    # 문서를 FAISS 벡터 DB에 저장할 때 1회 사용
    "text-embedding-3-small": {
        "tokens": 0.02 / 1_000_000,  # $0.02 / 1M tokens → 토큰당 $0.00000002
    },
}

# USD → KRW 환율 (참고용 고정값, 실제 환율과 다를 수 있음)
# 정확한 비용 확인은 OpenAI 대시보드(platform.openai.com/usage)를 이용하세요
USD_TO_KRW = 1_380


def _init_cost_state():
    """
    Streamlit 세션 상태(session_state)에 비용 추적 변수가 없으면 초기화합니다.

    Streamlit은 버튼 클릭 등 UI 이벤트 발생 시 스크립트 전체를 재실행합니다.
    session_state에 저장된 값은 재실행 사이에도 유지되므로,
    누적 비용을 잃지 않고 관리할 수 있습니다.
    """
    if "cost_log" not in st.session_state or st.session_state.cost_log is None:
        st.session_state.cost_log = []      # 각 API 호출 기록을 담는 리스트 (딕셔너리 항목들)
    if "total_cost_usd" not in st.session_state or st.session_state.total_cost_usd is None:
        st.session_state.total_cost_usd = 0.0   # 세션 전체 누적 비용 (USD)


def record_cost(label: str, model: str, **kwargs):
    """
    API 호출 한 건의 비용을 계산하여 session_state에 기록합니다.

    Args:
        label:  사용자에게 보여줄 호출 설명 (예: "핵심 정리 생성", "퀴즈 생성")
        model:  사용된 OpenAI 모델 이름 (PRICING 딕셔너리의 키와 일치해야 함)
        **kwargs:
            - input_tokens, output_tokens: GPT-4o 호출 시 사용 (토큰 수)
            - chars:                       TTS-1 호출 시 사용 (글자 수)
            - minutes:                     Whisper-1 호출 시 사용 (오디오 재생 시간(분))
            - tokens:                      임베딩 모델 호출 시 사용 (토큰 수)

    동작 흐름:
    1. 모델 종류에 따라 알맞은 요금 공식으로 비용(cost) 계산
    2. 계산된 비용과 사용량 정보를 딕셔너리로 만들어 cost_log 리스트에 추가
    3. total_cost_usd에 누적
    """
    _init_cost_state()  # 초기화 보장
    cost = 0.0
    detail = ""  # 사이드바에 표시할 사용량 설명 문자열

    if model == "gpt-4o":
        # GPT-4o: 입력 토큰과 출력 토큰을 각각 다른 단가로 계산 후 합산
        inp = kwargs.get("input_tokens", 0)
        out = kwargs.get("output_tokens", 0)
        cost = inp * PRICING["gpt-4o"]["input"] + out * PRICING["gpt-4o"]["output"]
        detail = f"입력 {inp:,} / 출력 {out:,} 토큰"  # 콤마 포맷(예: 1,234)으로 가독성 향상

    elif model == "tts-1":
        # TTS-1: 실제 생성된 오디오의 원본 텍스트 글자 수로 계산
        chars = kwargs.get("chars", 0)
        cost = chars * PRICING["tts-1"]["chars"]
        detail = f"{chars:,} 글자"

    elif model == "whisper-1":
        # Whisper-1: 오디오 재생 시간(분)으로 계산
        # 실제 재생 시간을 구하기 어려운 경우 파일 크기로 추정 (128kbps 기준)
        minutes = kwargs.get("minutes", 0)
        cost = minutes * PRICING["whisper-1"]["per_min"]
        detail = f"약 {minutes:.1f}분"

    elif model == "text-embedding-3-small":
        # Embedding: 벡터화된 총 토큰 수로 계산 (1토큰 ≈ 4글자 추정)
        tokens = kwargs.get("tokens", 0)
        cost = tokens * PRICING["text-embedding-3-small"]["tokens"]
        detail = f"{tokens:,} 토큰"

    # 기록 추가 및 누적
    st.session_state.cost_log.append({
        "label":    label,     # 호출 설명 (예: "플래시카드 생성")
        "model":    model,     # 모델명 (예: "gpt-4o")
        "detail":   detail,    # 사용량 요약 (예: "입력 3,200 / 출력 780 토큰")
        "cost_usd": cost,      # 해당 호출의 비용 (USD)
    })
    st.session_state.total_cost_usd += cost  # 세션 전체 누적 비용에 더하기


def render_cost_panel():
    """
    사이드바 하단에 API 비용 현황 패널을 렌더링합니다.

    표시 내용:
    - 세션 누적 비용 (USD + 원화 환산)
    - 접기/펼치기 가능한 호출 내역 상세 로그
    - 비용 기록 초기화 버튼

    ※ 이 함수는 사이드바 with 블록 안에서 호출되어야 합니다.
    """
    _init_cost_state()
    total = st.session_state.total_cost_usd  # 현재까지 누적된 총 비용
    log   = st.session_state.cost_log        # 호출 기록 리스트

    st.sidebar.markdown("---")  # 구분선
    st.sidebar.markdown("### 💸 API 사용 비용")

    if not log:
        # 아직 API 호출이 없으면 안내 메시지만 표시
        st.sidebar.caption("아직 API 호출 없음")
        return

    # --- 누적 비용 메트릭 카드 ---
    krw = total * USD_TO_KRW  # USD → KRW 환산
    st.sidebar.metric(
        label="누적 비용 (USD)",
        value=f"${total:.4f}",           # 소수점 4자리까지 표시 (예: $0.0342)
        delta=f"≈ ₩{krw:,.0f}",          # 콤마 포함 원화 (예: ≈ ₩47,196)
        delta_color="off",               # delta를 초록/빨강으로 표시하지 않음 (중립)
    )

    # --- 호출 내역 상세 로그 (접기/펼치기) ---
    with st.sidebar.expander("📋 호출 내역 보기"):
        # reversed(): 최신 호출이 맨 위에 오도록 역순으로 표시
        for i, entry in enumerate(reversed(log), 1):
            st.markdown(
                f"**{entry['label']}**  \n"        # 호출 설명 (굵게)
                f"`{entry['model']}` · {entry['detail']}  \n"  # 모델명 + 사용량
                f"→ `${entry['cost_usd']:.5f}`"    # 해당 호출 비용 (소수점 5자리)
            )
            if i < len(log):
                st.divider()  # 항목 사이 구분선 (마지막 항목 이후에는 생략)

    # --- 초기화 버튼 ---
    if st.sidebar.button("🗑️ 비용 기록 초기화"):
        # 버튼 클릭 시 비용 기록을 모두 지우고 화면을 새로고침
        st.session_state.cost_log = []
        st.session_state.total_cost_usd = 0.0
        st.rerun()  # Streamlit 전체 스크립트를 즉시 재실행하여 UI 반영


# ==============================================================================
# [9] LLM 클라이언트 캐싱 (성능 최적화)
# ==============================================================================
# Streamlit은 버튼 클릭마다 전체 스크립트를 재실행하므로,
# 캐싱 없이는 API 클라이언트가 매번 새로 생성됩니다.
# @st.cache_resource는 동일한 인자로 호출된 함수의 반환값을 메모리에 캐시하여 재사용합니다.
# (Fix #4: LLM 클라이언트 매번 재생성 문제 해결)

@st.cache_resource
def get_cached_chat_llm(api_key: str, temperature: float):
    """
    ChatOpenAI(GPT-4o) 인스턴스를 생성하고 캐시합니다.

    Args:
        api_key:     OpenAI API 키 (캐시 키의 일부로 사용)
        temperature: 생성 다양성 조절 (0.0=결정적, 1.0=창의적)
                     요약/시각화: 0.0~0.3 (일관성 중시)
                     플래시카드/퀴즈: 0.5 (약간의 다양성 허용)

    Returns:
        ChatOpenAI 인스턴스 (LangChain 래퍼)

    캐시 키: api_key + temperature 조합
    → 같은 키와 temperature로 호출하면 기존 객체를 재사용
    """
    return ChatOpenAI(model="gpt-4o", temperature=temperature, openai_api_key=api_key)


@st.cache_resource
def get_cached_openai_client(api_key: str):
    """
    OpenAI 공식 SDK 클라이언트를 생성하고 캐시합니다.
    LangChain을 거치지 않는 직접 API 호출
    (Whisper 전사, TTS, Vision, 플래시카드/퀴즈 JSON 생성)에 사용됩니다.

    Args:
        api_key: OpenAI API 키

    Returns:
        OpenAI 클라이언트 인스턴스
    """
    return OpenAI(api_key=api_key)


# ==============================================================================
# [10] 오디오 전처리 유틸리티 함수 (모듈 최상위 정의)
# ==============================================================================
# 원래 이 두 함수는 오디오 버튼 클릭 핸들러 내부에 중첩 정의되어 있었습니다.
# 버튼을 누를 때마다 함수가 재정의되어 비효율적이고 테스트하기 어려웠습니다.
# 모듈 최상위로 이동하여 재사용성과 가독성을 높였습니다. (Fix #3)

def clean_markdown_for_speech(text: str) -> str:
    """
    마크다운으로 작성된 요약문에서 TTS(음성 합성)에 불필요한 기호를 제거합니다.

    제거 대상:
    - 헤더 기호: ### 핵심 내용 → 핵심 내용
    - 볼드체:    **중요한 내용** → 중요한 내용
    - 밑줄체:    __중요한 내용__ → 중요한 내용
    - 연속 줄바꿈: \n\n → 공백 한 칸 (음성에서 어색한 침묵 방지)

    Args:
        text: 마크다운 형식의 원본 텍스트

    Returns:
        기호가 제거된 평문 텍스트 (앞뒤 공백 제거됨)
    """
    text = re.sub(r'#+\s?', '', text)   # '#', '##', '###' 등 헤더 기호 제거
    text = re.sub(r'\*\*|__', '', text) # '**'(볼드) 및 '__'(밑줄) 기호 제거
    text = re.sub(r'\n+', ' ', text)    # 하나 이상의 줄바꿈을 공백 한 칸으로 교체
    return text.strip()                 # 양쪽 끝 불필요한 공백 제거


def extract_all_core_parts(text: str, ui_text: dict) -> str:
    """
    전체 요약 텍스트에서 "핵심 내용 요약" 섹션만 추출하여 반환합니다.

    TTS 오디오를 생성할 때 표나 용어 정리 섹션은 읽기 어색하므로,
    '핵심 내용 요약' 헤더(h_bullet)와 '상세 요약 표' 헤더(h_table) 사이의
    내용만 골라내어 자연스러운 오디오 스크립트를 만듭니다.

    예시:
        입력 텍스트에 "1. 핵심 내용 요약 ... 2. 상세 요약 표" 패턴이 여러 개 있으면
        모두 추출하여 하나로 합칩니다 (강의 여러 개를 한 번에 브리핑).

    Args:
        text:     gen_summary()가 반환한 전체 마크다운 요약 텍스트
        ui_text:  현재 언어 팩 딕셔너리 (h_bullet, h_table 키 사용)

    Returns:
        추출되고 정제된 평문 텍스트
        (패턴 매칭 실패 시 원본 텍스트 앞 1000자를 정제하여 반환)
    """
    start_marker = ui_text['h_bullet']  # 시작 마커: "1. 핵심 내용 요약" 또는 "1. Key Highlights"
    end_marker   = ui_text['h_table']   # 종료 마커: "2. 상세 요약 표" 또는 "2. Detailed Summary Table"

    # re.escape(): 마커 문자열에 정규표현식 특수문자가 있어도 안전하게 처리
    # re.DOTALL: '.'이 줄바꿈 문자도 포함하여 매칭 (멀티라인 섹션 추출 필수)
    # (.*?): 비탐욕(non-greedy) 매칭으로 가장 짧은 구간만 캡처
    pattern = f"{re.escape(start_marker)}(.*?){re.escape(end_marker)}"
    matches = re.findall(pattern, text, re.DOTALL)

    if matches:
        # 여러 섹션(강의)이 있으면 공백으로 이어 붙여 하나의 스크립트로 만듦
        combined_text = " ".join(matches)
        return clean_markdown_for_speech(combined_text)
    else:
        # 패턴 매칭 실패 (전체 모드가 아닌 경우 등): 원본 앞 1000자를 대신 사용
        return clean_markdown_for_speech(text[:1000])


# ==============================================================================
# [11] 핵심 로직: 파일 텍스트 추출
# ==============================================================================

def extract_text(file, ext: str, path: str, key: str) -> str:
    """
    파일 형식에 맞는 방법으로 텍스트를 추출하여 반환합니다.

    지원 형식 및 처리 방식:
    - PDF (.pdf):       LangChain PyPDFLoader → 페이지별 텍스트 추출 후 합치기
    - Word (.docx/.doc): LangChain Docx2txtLoader → 문서 전체 텍스트 추출
    - PPT (.pptx/.ppt): python-pptx → 슬라이드별 도형(shape)에서 텍스트 추출
    - 이미지 (.jpg/.png/.jpeg): GPT-4o Vision API → Base64 인코딩 후 이미지에서 텍스트 인식
    - 오디오 (.mp3/.wav/.m4a): Whisper API → 음성을 텍스트로 전사
    - 영상 (.mp4/.avi/.mov):   MoviePy로 오디오 트랙 추출 → Whisper API 전사

    Args:
        file: Streamlit UploadedFile 객체 (file.name으로 파일명 접근)
        ext:  파일 확장자 소문자 (예: ".pdf", ".mp4")
        path: 임시 저장된 파일의 전체 경로
        key:  OpenAI API 키

    Returns:
        추출된 텍스트 문자열 (실패 시 에러 메시지 또는 빈 문자열)
    """

    # --- PDF 처리 ---
    if ext == ".pdf":
        # PyPDFLoader는 PDF를 페이지 단위 Document 리스트로 반환
        # page_content를 모두 이어붙여 하나의 문자열로 만듦
        return "".join([p.page_content for p in PyPDFLoader(path).load()])

    # --- Word 문서 처리 ---
    elif ext in [".docx", ".doc"]:
        # Docx2txtLoader도 Document 리스트 반환 (대부분 1개 항목)
        return "".join([p.page_content for p in Docx2txtLoader(path).load()])

    # --- PowerPoint 처리 ---
    elif ext in [".pptx", ".ppt"]:
        prs = Presentation(path)  # PPT 파일 파싱
        # 모든 슬라이드(prs.slides) → 모든 도형(sl.shapes) → 텍스트가 있는 도형만 선택
        # hasattr(s, "text"): 텍스트 박스가 아닌 이미지나 선 도형은 건너뜀
        return "\n".join([s.text for sl in prs.slides for s in sl.shapes if hasattr(s, "text")])

    # --- 이미지 처리 (GPT-4o Vision) ---
    elif ext in [".jpg", ".png", ".jpeg"]:
        # OpenAI SDK를 직접 사용 (LangChain의 Vision 지원이 불안정하므로)
        client = OpenAI(api_key=key)
        # 이미지를 Base64로 인코딩: OpenAI Vision API는 URL 또는 Base64만 허용
        with open(path, "rb") as f:
            enc = base64.b64encode(f.read()).decode('utf-8')
        # data URL 형식으로 이미지 전달: "data:image/jpeg;base64,{실제 데이터}"
        res = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": [
                {"type": "text",      "text": "Extract all text visible in this slide/image precisely."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{enc}"}}
            ]}]
        )
        # 이미지 처리에 사용된 토큰 수를 비용 시스템에 기록
        record_cost("이미지 텍스트 추출", "gpt-4o",
                    input_tokens=res.usage.prompt_tokens,
                    output_tokens=res.usage.completion_tokens)
        # 출처 정보를 텍스트 앞에 태그로 붙여 반환 (나중에 AI가 출처를 파악하는 데 도움)
        return f"[Image Source: {file.name}] " + res.choices[0].message.content

    # --- 오디오 처리 (Whisper) ---
    elif ext in [".mp3", ".wav", ".m4a"]:
        client = OpenAI(api_key=key)
        # Whisper API는 재생 시간을 반환하지 않으므로, 파일 크기로 재생 시간을 추정
        # 추정 공식: (파일 바이트 수) / (128kbps → 초당 바이트 수) / 60 = 분
        # 128kbps = 128 * 1024 bits/s = 128 * 1024 / 8 bytes/s ≈ 16,384 bytes/s
        with open(path, "rb") as f:
            audio_bytes = f.read()
        estimated_minutes = len(audio_bytes) / (128 * 1024 / 8) / 60
        # Whisper API 호출: 파일을 그대로 전송하여 텍스트로 변환
        with open(path, "rb") as f:
            txt = client.audio.transcriptions.create(model="whisper-1", file=f).text
        record_cost("음성 전사 (Whisper)", "whisper-1", minutes=estimated_minutes)
        return f"[Audio Source: {file.name}] " + txt

    # --- 영상 처리 (MoviePy → Whisper) ---
    elif ext in [".mp4", ".avi", ".mov"]:
        # MoviePy가 설치되어 있지 않으면 처리 불가
        if VideoFileClip is None:
            return "Error: MoviePy missing."

        # 영상에서 오디오만 추출하여 임시 MP3 파일로 저장
        audio_path = path + "_temp.mp3"  # 원본 경로에 "_temp.mp3"를 붙여 임시 파일명 생성
        vid = None  # finally 블록에서 vid.close()를 안전하게 호출하기 위해 미리 None으로 초기화
        try:
            vid = VideoFileClip(path)  # 영상 파일 로드 (메모리 + 파일 핸들 점유)
            # 오디오 트랙을 MP3로 추출 (logger=None: 진행 로그 출력 억제)
            vid.audio.write_audiofile(audio_path, logger=None)
            # 추출된 MP3를 Whisper로 전사
            client = OpenAI(api_key=key)
            with open(audio_path, "rb") as f:
                txt = client.audio.transcriptions.create(model="whisper-1", file=f).text
            # 영상의 실제 재생 시간(초)을 분으로 변환하여 비용 기록 (추정값보다 정확)
            estimated_minutes = vid.duration / 60
            record_cost("영상 전사 (Whisper)", "whisper-1", minutes=estimated_minutes)
            return f"[Video Source: {file.name}] " + txt
        except Exception as e:
            return str(e)  # 처리 실패 시 에러 메시지를 텍스트로 반환
        finally:
            # finally 블록은 try/except 완료 후 무조건 실행됨
            # Fix #2: vid.close() 누락으로 인한 파일 핸들 누수 방지
            # - Windows: close 없이 파일 삭제 시 "파일이 다른 프로세스에서 사용 중" 오류 발생
            # - 장기 실행 서버: close 없으면 메모리 누수가 점진적으로 증가
            if vid is not None:
                vid.close()
            # 임시 오디오 파일 삭제 (존재할 때만)
            if os.path.exists(audio_path):
                os.remove(audio_path)

    # 지원하지 않는 확장자: 빈 문자열 반환 → 해당 파일은 DB에 추가되지 않음
    return ""


# ==============================================================================
# [12] 핵심 로직: 지식베이스(벡터 DB) 구축
# ==============================================================================

def build_knowledge_base(lec_files, prob_files, key: str, ui_text: dict):
    """
    업로드된 파일들을 읽고, 벡터 DB(FAISS)를 구축하는 공개 인터페이스입니다.

    실제 무거운 작업은 _build_knowledge_base_cached()에 위임하며,
    이 함수는 파일 내용 해시를 계산하여 동일한 파일 조합에 대한
    캐시 히트 여부를 결정하는 역할만 합니다. (Fix #6)

    동작 원리 (캐싱):
    - 각 파일의 내용(bytes)을 MD5 해시로 변환
    - 모든 파일의 해시를 튜플로 묶어 캐시 키 생성
    - 동일한 파일 조합이라면 API 재호출 없이 캐시된 DB 반환
    - 파일이 하나라도 변경되면 새로운 DB를 생성

    ※ st.cache_resource는 함수 인자가 hashable해야 하므로 list 대신 tuple 사용

    Args:
        lec_files:  강의 자료 파일 목록 (Streamlit UploadedFile 리스트)
        prob_files:  연습 문제 파일 목록
        key:        OpenAI API 키
        ui_text:    현재 언어 팩 딕셔너리

    Returns:
        FAISS 벡터 DB 객체 (실패 시 None)
    """
    import hashlib  # 파일 내용의 MD5 해시를 계산하기 위한 표준 라이브러리

    def _file_hash(f) -> str:
        """Streamlit UploadedFile의 내용을 MD5 해시 문자열로 변환"""
        # f.getvalue(): 파일 전체 내용을 bytes로 반환
        return hashlib.md5(f.getvalue()).hexdigest()

    # 강의 자료 + 연습 문제 파일을 합쳐서 각각 해시 계산
    # tuple로 변환: list는 hashable하지 않아 @st.cache_resource의 캐시 키로 사용 불가
    all_hashes = tuple(
        _file_hash(f) for f in (lec_files or []) + (prob_files or [])
    )
    # 해시 튜플을 캐시 키로 넘겨서 실제 DB 생성 함수 호출
    return _build_knowledge_base_cached(all_hashes, lec_files, prob_files, key, ui_text)


@st.cache_resource(show_spinner=False)
def _build_knowledge_base_cached(file_hashes: tuple, lec_files, prob_files, key: str, ui_text: dict):
    """
    실제 DB 생성 로직을 수행하는 캐시된 내부 함수입니다.

    @st.cache_resource: file_hashes(파일 내용 해시 튜플)가 동일하면
    이 함수를 실행하지 않고 이전에 반환한 FAISS DB 객체를 그대로 반환합니다.
    → 동일한 파일을 다시 업로드해도 임베딩 API를 재호출하지 않아 비용 절감.

    처리 파이프라인:
    1. 각 파일을 임시 디스크에 저장 (Streamlit UploadedFile은 메모리 객체)
    2. extract_text()로 형식별 텍스트 추출
    3. 소스 타입(강의 자료/연습 문제)을 태그로 붙여 Document 생성
    4. RecursiveCharacterTextSplitter로 청크 분할 (1500자 단위, 200자 오버랩)
    5. OpenAIEmbeddings로 청크를 벡터로 변환
    6. FAISS.from_documents()로 벡터 DB 구축

    Args:
        file_hashes: 파일 내용 해시 튜플 (캐시 키로만 사용, 함수 내에서 직접 사용 안 함)
        (나머지 인자는 build_knowledge_base()와 동일)
    """
    docs = []  # 처리된 Document 객체를 담을 리스트
    # st.status(): 처리 진행 상황을 사이드바에 실시간으로 표시하는 UI 컴포넌트
    status = st.status(ui_text["msg_proc"], expanded=True)

    def process_files(file_list, source_type: str):
        """
        파일 리스트를 순서대로 처리하여 docs 리스트에 Document를 추가하는 내부 함수.

        각 파일은:
        1. 임시 파일로 저장 (extract_text가 파일 경로를 필요로 함)
        2. 텍스트 추출
        3. [소스 타입] 태그를 앞에 붙여 Document로 래핑
        4. 처리 완료 후 임시 파일 삭제

        Args:
            file_list:   처리할 Streamlit UploadedFile 리스트
            source_type: "Lecture Material" 또는 "Practice Problem"
                         → AI가 강의 내용과 연습 문제를 구분하여 처리하는 데 활용
        """
        for f in file_list:
            ext = os.path.splitext(f.name)[1].lower()  # 확장자 추출 및 소문자 변환
            status.write(f"{ui_text['msg_ingest']} [{source_type}] {f.name}")  # 진행 상황 표시

            # NamedTemporaryFile: 이름이 있는 임시 파일 생성
            # delete=False: close() 후에도 파일이 자동 삭제되지 않도록 설정 (수동 삭제 예정)
            # suffix=ext: extract_text()가 확장자를 통해 파일 형식을 판별하므로 필수
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(f.getvalue())  # 업로드된 파일 내용을 디스크에 저장
                tmp_path = tmp.name      # 임시 파일 경로 저장

            try:
                content = extract_text(f, ext, tmp_path, key)  # 형식별 텍스트 추출
                if content:
                    # [소스 타입] 태그를 텍스트 앞에 명시
                    # → AI 프롬프트에서 강의 자료와 연습 문제를 구분하여 다르게 처리
                    # 예: "[Lecture Material] \n신경망이란 ..." 또는 "[Practice Problem] \n문제 1. ..."
                    tagged_content = f"[{source_type}] \n{content}"
                    docs.append(Document(
                        page_content=tagged_content,
                        metadata={"source": f.name, "type": source_type}
                        # metadata: 나중에 다이어그램 생성 시 파일명을 추출하는 데 활용
                    ))
            except Exception as e:
                st.error(f"{ui_text['msg_err_file']} {e}")  # 처리 실패 시 에러 메시지 표시
            finally:
                os.remove(tmp_path)  # 성공/실패 여부와 관계없이 임시 파일 삭제

    # 강의 자료와 연습 문제를 각각 다른 source_type으로 처리
    if lec_files:  process_files(lec_files,  "Lecture Material")
    if prob_files: process_files(prob_files, "Practice Problem")

    # 유효한 Document가 하나도 없으면 실패 처리
    if not docs:
        status.update(label=ui_text["msg_nodata"], state="error")
        return None

    # --- 텍스트 청크 분할 ---
    # chunk_size=1500: 한 청크의 최대 글자 수
    # chunk_overlap=200: 인접 청크 간 겹치는 글자 수 (청크 경계에서 문맥이 잘리지 않도록)
    # RecursiveCharacterTextSplitter는 ["\n\n", "\n", " ", ""] 순으로 분할 지점을 시도하여
    # 가능한 한 단락/문장/단어 경계에서 자름
    splitter = RecursiveCharacterTextSplitter(chunk_size=1500, chunk_overlap=200)
    splits = splitter.split_documents(docs)  # Document 리스트 → 더 작은 청크 Document 리스트

    # --- 임베딩 비용 추정 및 기록 ---
    # OpenAIEmbeddings는 내부적으로 토큰 사용량을 반환하지 않으므로 글자 수로 추정
    # 영어: 1토큰 ≈ 4글자, 한국어: 1토큰 ≈ 2~3글자 → 평균 4글자로 보수적 추정
    total_chars = sum(len(d.page_content) for d in splits)
    estimated_embed_tokens = total_chars // 4  # 글자 수 / 4 = 추정 토큰 수

    # --- FAISS 벡터 DB 생성 ---
    # OpenAIEmbeddings: 각 청크를 OpenAI 임베딩 API로 벡터로 변환
    # text-embedding-3-small이 내부적으로 사용됨 (LangChain 기본값)
    # FAISS.from_documents(): 모든 벡터를 FAISS 인덱스에 저장
    db = FAISS.from_documents(splits, OpenAIEmbeddings(openai_api_key=key))
    record_cost("임베딩 (Embeddings)", "text-embedding-3-small", tokens=estimated_embed_tokens)

    # 처리 완료 표시 (st.status 컴포넌트를 접힌 상태로 업데이트)
    status.update(label=ui_text["msg_done"], state="complete", expanded=False)
    return db  # 완성된 FAISS 벡터 DB 반환


# ==============================================================================
# [13] 핵심 로직: RAG QA 체인 생성
# ==============================================================================

def get_rag_chain(db, key: str, target_lang: str):
    """
    AI 도우미 탭에서 사용할 RAG (Retrieval-Augmented Generation) 체인을 생성합니다.

    RAG 체인 동작 흐름:
    1. 사용자 질문 수신
    2. FAISS DB에서 관련 청크 검색 (MMR 방식)
    3. 검색된 청크를 컨텍스트로 삼아 GPT-4o에 답변 생성 요청
    4. 답변 반환

    검색 방식 (MMR - Maximal Marginal Relevance):
    - 단순 유사도 검색과 달리 다양성을 함께 고려
    - 중복 내용의 청크보다 서로 다른 측면을 커버하는 청크 조합을 선호
    - lambda_mult: 0.0(다양성 최대) ~ 1.0(유사도 최대), 0.6은 균형점

    프롬프트 구조:
    - SOURCE IDENTIFICATION: 강의 자료와 연습 문제를 구분하여 다르게 처리
    - INSTRUCTIONS: 개념 설명 모드 / 시험 전략 모드 구분
    - STRICT RULES: 컨텍스트 외 정보 사용 금지, 출력 언어 지정

    Args:
        db:          FAISS 벡터 DB 객체
        key:         OpenAI API 키
        target_lang: 답변 언어 ("Korean" 또는 "English")

    Returns:
        LangChain RetrievalQA 체인 객체 (chain.invoke({"query": ...})로 호출)
    """
    # 캐시된 LLM 클라이언트 사용 (temperature=0.2: 일관성 있는 답변 선호)
    llm = get_cached_chat_llm(key, temperature=0.2)

    # --- 시스템 프롬프트 템플릿 ---
    # f-string으로 target_lang을 삽입하여 출력 언어를 지정
    # {context}와 {question}은 LangChain이 실행 시 자동으로 채워넣는 자리표시자
    # {{context}}처럼 이중 중괄호를 쓰는 이유: f-string에서 { }는 변수이므로
    # 리터럴 { }를 표현하려면 {{ }}로 이스케이프해야 함
    template = f"""
    You are an intelligent AI Teaching Assistant and Exam Strategist.
    
    *** SOURCE IDENTIFICATION ***
    - Text starting with `[Lecture Material]` is conceptually explanatory.
    - Text starting with `[Practice Problem]` contains actual exam/quiz questions.

    *** INSTRUCTIONS ***
    1. **Concept Explainer**: If asked about concepts, prioritize `[Lecture Material]`.
    2. **Exam Strategist**: If asked about "exam style", "preparation", or "types of problems":
       - Look strictly at the content labeled `[Practice Problem]`.
       - Analyze the format (Multiple choice? Essay? Calculation?) and difficulty.
       - Provide a strategy based on those specific patterns.
       - If no `[Practice Problem]` is found, state that you need practice files to analyze the exam style.
    
    *** STRICT RULES ***
    - Answer ONLY using the provided [Context].
    - Output Language: **{target_lang}**.
    
    [Context]:
    {{context}}
    
    [User Question]:
    {{question}}
    """

    # --- MMR 검색기 설정 ---
    retriever = db.as_retriever(
        search_type="mmr",       # Maximal Marginal Relevance: 다양성을 고려한 검색
        search_kwargs={
            'k': 20,             # 최종 반환할 청크 수 (많을수록 컨텍스트 풍부)
            'fetch_k': 50,       # MMR 계산을 위해 후보로 가져올 청크 수 (k보다 커야 함)
            'lambda_mult': 0.6   # 유사도(0.6)와 다양성(0.4)의 가중치 균형
        }
    )

    # --- RetrievalQA 체인 조립 ---
    return RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever,
        # chain_type_kwargs: 프롬프트 템플릿을 체인에 주입
        # input_variables는 LangChain이 자동으로 채우는 변수명 목록
        chain_type_kwargs={
            "prompt": PromptTemplate(
                template=template,
                input_variables=["context", "question"]
            )
        }
    )


# ==============================================================================
# [14] 생성 함수 공통 유틸리티
# ==============================================================================

def get_scope(topic: str) -> str:
    """topic이 비어있으면 전체 범위, 아니면 특정 토픽 문자열 반환 (프롬프트 삽입용)"""
    return "the ENTIRE provided material" if not topic or not topic.strip() else f"the topic '{topic}'"


def clean_json(text: str) -> str:
    """
    LLM이 반환한 텍스트에서 순수 JSON 배열만 추출하여 반환합니다.

    LLM은 종종 JSON 앞뒤에 마크다운 코드블록(```json ... ```)이나
    설명 문자열을 붙여 반환하므로, 이를 제거하고 [ ... ] 부분만 파싱합니다.

    처리 순서:
    1. 앞뒤 공백 제거
    2. ```json 또는 ``` 마크다운 코드블록 제거
    3. 정규표현식으로 JSON 배열 [ ... ] 추출
    4. 후행 콤마 제거 (,] → ]) - 일부 LLM이 마지막 항목 뒤에 콤마를 붙이는 버그 대응
    """
    text = text.strip()
    # re.IGNORECASE: ```JSON, ```Json 등 대소문자 무관하게 제거
    text = re.sub(r"```(json)?", "", text, flags=re.IGNORECASE).replace("```", "")
    # re.DOTALL: '.'이 줄바꿈도 포함하므로 여러 줄에 걸친 JSON 배열 추출 가능
    match = re.search(r"(\[.*\])", text, re.DOTALL)
    if match:
        text = match.group(1)  # 첫 번째 캡처 그룹: [ ... ] 전체
    text = re.sub(r",\s*\]", "]", text)  # trailing comma 제거 (예: [1, 2, 3,] → [1, 2, 3])
    return text


def clean_dot_code(text: str) -> str:
    """
    LLM이 반환한 텍스트에서 유효한 Graphviz DOT 코드만 추출합니다.

    LLM은 DOT 코드 앞뒤에 설명이나 마크다운 코드블록을 붙이는 경향이 있으므로,
    "digraph" 키워드부터 마지막 "}" 사이의 내용만 잘라냅니다.

    처리 순서:
    1. ```dot 마크다운 코드블록 제거
    2. "digraph" 키워드 위치 찾기
    3. 첫 번째 "{" 위치 찾기
    4. 마지막 "}" 위치 찾기
    5. "digraph ~ }" 범위의 텍스트만 반환
    """
    text = text.strip()
    text = re.sub(r"```(dot)?", "", text).replace("```", "")  # 마크다운 코드블록 제거
    start_idx = text.find("digraph")   # "digraph G {" 시작 위치
    if start_idx == -1: return text    # "digraph" 없으면 원본 반환 (에러 케이스)
    open_brace  = text.find("{", start_idx)  # 첫 번째 { 위치
    if open_brace == -1: return text
    close_brace = text.rfind("}")     # 마지막 } 위치 (rfind: 오른쪽부터 검색)
    if close_brace == -1: return text
    return text[start_idx : close_brace + 1]  # "digraph ... }" 슬라이싱


# ==============================================================================
# [15] 생성 함수: 핵심 정리 (요약)
# ==============================================================================

def gen_summary(db, api_key: str, topic: str, ui_text: dict) -> str:
    """
    벡터 DB에서 관련 내용을 검색하여 핵심 정리(요약)를 생성합니다.

    두 가지 모드:
    1. 전체 모드 (topic 비어있거나 "전체" 등):
       - DB에서 강의 목차/제목 관련 80개 청크 검색
       - 강의별로 간략한 개요 + 핵심 시험 개념 목록 생성
       - 표(테이블) 없는 경량화 포맷 (전체 강의를 한 출력에 담기 위해)

    2. 특정 토픽 모드 (예: "신경망"):
       - 해당 주제 관련 15개 청크 집중 검색
       - 핵심 포인트 + 상세 요약 표 + 용어 정리 표 생성
       - 풍부한 상세 포맷

    Args:
        db:       FAISS 벡터 DB
        api_key:  OpenAI API 키
        topic:    사용자 입력 주제 (빈 문자열 또는 전체 모드 키워드 가능)
        ui_text:  현재 언어 팩

    Returns:
        마크다운 형식의 요약 문자열 (에러 시 에러 메시지 문자열)
    """
    # 전체 모드 여부 판별 및 LLM에 전달할 범위 설명 문자열 설정
    is_all_mode = check_is_all_mode(topic)
    if is_all_mode:
        scope_text = "the ENTIRE provided material (All lectures)"
    else:
        scope_text = f"the specific topic '{topic}'"

    lang = ui_text["target_lang"]  # 출력 언어 ("Korean" 또는 "English")

    # --- 모드별 설정 ---
    if is_all_mode:
        # [전체 모드] 전체 강의 개요를 효율적으로 커버하기 위한 설정
        search_query = "Table of contents, Lecture titles, Course outline, Key concepts summary"
        # k_val=80: 많은 청크를 검색하여 전체 강의를 빠짐없이 커버
        k_val = 80

        mode_instruction = f"""
        - **Goal**: Create a **"Master Course Outline"** that lists **EVERY** detected file or lecture.
        - **Constraint**: Keep descriptions concise to ensure ALL lectures are covered within the output limit.
        - **Coverage**: It is critical to list **ALL** lectures/files found in the text. Do not stop after the first few.
        - **Format**: For each lecture, provide a brief summary and a list of key exam concepts.
        """

        guidelines = f"""
        1. **Context-Based**: Answer ONLY based on the provided [Context].
        2. **Completeness (CRITICAL)**: 
           - You MUST iterate through **ALL** detected files/lectures.
           - Do not skip the later lectures. 
        3. **Terminology Integrity (STRICT)**: 
           - Even in the concept list, terms must be a **VERBATIM COPY** from the source.
           - **DO NOT TRANSLATE THE TERM ITSELF.**
        """

        # 표 없는 경량 포맷: 강의가 많을 때 하나의 출력 안에 전체를 담을 수 있도록 최소화
        format_instruction = f"""
        **[INSTRUCTION: Repeat the block below for EVERY detected File/Lecture]**

        ## 📂 [Insert File or Lecture Name]
        
        ### 📖 {ui_text['h_bullet']} (Lecture Overview)
        - (Summarize the main theme of this lecture in 2-3 sentences in {lang}.)
        
        ### 🔑 Key Exam Concepts
        - **(Concept 1)**: (Short definition/Core logic in {lang})
        - **(Concept 2)**: (Short definition/Core logic in {lang})
        - **(Concept 3)**: (Short definition/Core logic in {lang})
        
        ---
        """

    else:
        # [특정 토픽 모드] 선택된 주제에 집중하여 심층 분석
        search_query = topic   # 사용자가 입력한 주제를 그대로 검색어로 사용
        k_val = 15  # 특정 주제 관련 청크만 집중적으로 가져옴

        mode_instruction = f"""
        - **Scope Focus**: Focus **DEEPLY and STRICTLY** on the concept of '{topic}'. Ignore unrelated sections.
        - **Terminology Strategy**: Select terms that are **semantically related** to '{topic}' (e.g., sub-concepts, components, algorithms).
        """

        guidelines = f"""
        1. **Context-Based**: Answer ONLY based on the provided [Context].
        2. **Comprehensive Coverage**: 
           - Do NOT limit the number of key points.
           - Extract **ALL** core concepts, definitions, formulas, and arguments.
           - Aim for high detail.
        3. **Terminology Integrity (STRICT)**: 
           - In the Terminology Table, the 'Term' column must be a **VERBATIM COPY** from the source text.
           - **DO NOT TRANSLATE THE TERM ITSELF.**
           - If the source uses English (e.g., "Backpropagation"), keep it "Backpropagation".
           - If the source uses Korean (e.g., "역전파"), keep it "역전파".
           - Only the Definition and Context columns should be in **{lang}**.
        """

        # 3가지 섹션 포맷: 핵심 정리 + 상세 표 + 용어 정리 표
        format_instruction = f"""
        ### {ui_text['h_bullet']}
        - (List **ALL** exam-relevant key points about '{topic}' in {lang}.)
        
        ### {ui_text['h_table']}
        | {ui_text['h_th'][0]} | {ui_text['h_th'][1]} |
        |---|---|
        | (Category in {lang}) | (Detailed explanation in {lang}) |
        
        ### {ui_text['h_term']}
        | {ui_text['h_th'][2]} | {ui_text['h_th'][3]} | {ui_text['h_th'][4]} |
        |---|---|---|
        | **(EXACT SOURCE TERM)** | (Definition in {lang}) | (Context/Relation in {lang}) |
        """

    # --- DB 유효성 검사 ---
    if db is None:
        return "Error: Database is not initialized."

    # --- 벡터 검색 ---
    # similarity_search(): 검색어와 코사인 유사도가 높은 청크 k_val개 반환
    docs = db.similarity_search(search_query, k=k_val)
    # 검색된 청크들의 텍스트를 줄바꿈으로 이어붙여 하나의 컨텍스트 문자열 생성
    context = "\n".join([d.page_content for d in docs])

    # --- 최종 프롬프트 조합 ---
    prompt = f"""
    You are an expert **Professor** and **Exam Preparation Tutor**.
    Analyze {scope_text} based STRICTLY on the provided context.
    
    *** MODE INSTRUCTION ***
    {mode_instruction}
    
    *** CRITICAL GUIDELINES ***
    {guidelines}

    *** OUTPUT FORMAT ***
    {format_instruction}

    [Context]:
    {context}
    """

    # --- LLM 호출 + 비용 추적 ---
    try:
        llm = get_cached_chat_llm(api_key, temperature=0.3)
        # get_openai_callback(): with 블록 내 LangChain LLM 호출의 토큰 사용량 자동 집계
        with get_openai_callback() as cb:
            response = llm.invoke(prompt)
        record_cost("핵심 정리 생성", "gpt-4o",
                    input_tokens=cb.prompt_tokens,      # 프롬프트(입력) 토큰 수
                    output_tokens=cb.completion_tokens) # 응답(출력) 토큰 수
        return response.content  # LangChain AIMessage 객체에서 텍스트만 추출
    except Exception as e:
        return f"Error during generation: {str(e)}"


# ==============================================================================
# [16] 생성 함수: 시각화 (마인드맵 / 스파이더 다이어그램)
# ==============================================================================

def gen_diagram_optimized(db, api_key: str, topic: str, viz_type: str, ui_text: dict) -> str:
    """
    강의 내용을 Graphviz DOT 형식의 다이어그램 코드로 생성합니다.

    두 가지 다이어그램 유형:
    - Mindmap:       계층형 트리 구조 (상위 → 하위 개념), 수평 방향 배치 (dot 엔진)
    - Spider Diagram: 방사형 구조 (중심 주제 → 관련 개념), 방사형 배치 (neato 엔진)

    두 가지 내용 모드:
    - 전체 모드: 모든 강의 목록을 노드로 표현 (Course Overview → 각 강의 → 핵심 키워드)
    - 특정 토픽: 하나의 개념을 세부 구성요소로 분해하여 표현

    DOT 코드 생성 전략:
    1. 적절한 검색 방식으로 관련 청크 수집
    2. 각 청크에 파일명(출처)을 명시하여 AI가 누락 없이 전체 강의를 인식하도록 함
    3. GPT-4o에 Graphviz DOT 코드 생성 요청 (엄격한 규칙: 엣지에 텍스트 금지 등)
    4. 반환된 텍스트에서 clean_dot_code()로 순수 DOT 코드만 추출

    Args:
        db:       FAISS 벡터 DB
        api_key:  OpenAI API 키
        topic:    다이어그램 주제 (비어있으면 전체 강의 구조)
        viz_type: "Mindmap" 또는 "Spider Diagram"
        ui_text:  현재 언어 팩

    Returns:
        Graphviz DOT 코드 문자열 (에러 시 에러 메시지가 담긴 최소 DOT 코드)
    """
    # 1. 전체/특정 토픽 모드 판별
    is_all_mode = check_is_all_mode(topic)

    # 2. 검색 전략 및 프롬프트 지침 설정
    if is_all_mode:
        # [전체 모드]: 강의 목차/목록을 찾기 위한 포괄적 검색어
        search_query = "Table of Contents, Course Syllabus, All Lecture Titles, All Chapter Titles, Lecture 1, Lecture 2, ..., Lecture N"
        search_type = "mmr"  # 다양한 강의를 골고루 포함하기 위해 다양성 보장 검색 사용

        # Fix #5: db.index.ntotal로 실제 청크 수를 확인하여 k와 fetch_k를 동적으로 조정
        # 기존 k=300, fetch_k=3000은 실제 DB 크기를 훨씬 초과하여 불필요한 연산 낭비
        total_docs = db.index.ntotal             # FAISS 인덱스에 저장된 실제 벡터(청크) 총 수
        k_val   = min(100, total_docs)           # 최대 100개, DB 크기 초과 방지
        fetch_k = min(500, total_docs)           # MMR 후보 풀, 역시 DB 크기 초과 방지

        root_node = "Course Overview"  # 다이어그램 최상위 루트 노드 레이블

        scope_instruction = """
        - **MODE**: Full Course Syllabus & Key Concepts.
        - **GOAL**: Visualize **EVERY SINGLE** Lecture/Chapter found in the files, and optionally attach 2-3 key concepts to each lecture.
        - **CRITICAL REQUIREMENT (NO OMISSION)**: 
            1. **Exhaustive List**: Look at the [Source File] names and context. You MUST create a node for EVERY lecture present (e.g., Lecture 1 to Lecture N). **DO NOT SKIP ANY LECTURE.**
            2. **Hierarchy**: Root -> Lecture Node (Level 1) -> Keyword Nodes (Level 2).
            3. **NO EDGE LABELS**: Edges must be plain lines. **Put all text INSIDE the Node.**
            4. **Logical Order**: Arrange nodes in the order.
        - **NAMING RULES (CRITICAL)**:
            1. Node Label: "Lec X: [Title]" (e.g., "Lec 2: Metals", "Lec 5: Composites").
               - BAD: "Lecture 2" -> "Metals" (Do not split).
               - GOOD: Root -> "Lec 2: Metals".
            2. Keyword Label: Use the exact term from the context (e.g., "Thermodynamics", "Stress-Strain Curve").
            3. **NO EDGE LABELS**: Edges must be plain lines. Text goes inside nodes.
        """

    else:
        # [특정 토픽 모드]: 해당 주제의 구조를 파악하기 위한 집중 검색
        search_query = f"Structure and details of '{topic}', sub-types, components, key features"
        search_type = "similarity"  # 특정 주제는 일반 유사도 검색으로 충분
        k_val   = 15    # 특정 주제 집중 검색은 소수의 고관련도 청크로 충분
        fetch_k = 0     # MMR 사용 안 함 (similarity 모드에서는 미사용)

        root_node = topic.strip()  # 입력한 주제명을 루트 노드 레이블로 사용

        scope_instruction = f"""
        - **MODE**: Structured Deep Dive.
        - **GOAL**: Visualize the **Structure** of '{topic}' concisely.
        - **STYLE**:
            1. Root ('{topic}') -> Sub-Components / Types (Level 1).
            2. Sub-Components -> Key Characteristics (Level 2).
            3. **Constraint**: Use short phrases in nodes (Max 5-8 words). Avoid long sentences.
            4. **NO EDGE LABELS**: Edges must be plain lines. Text goes inside nodes.
        """

    # 3. 벡터 DB 검색
    try:
        if search_type == "mmr":
            # max_marginal_relevance_search(): fetch_k개 후보 중 다양성을 고려하여 k개 선별
            docs = db.max_marginal_relevance_search(search_query, k=k_val, fetch_k=fetch_k)
        else:
            # similarity_search(): 코사인 유사도 기준 상위 k개 반환
            docs = db.similarity_search(search_query, k=k_val)
    except Exception as e:
        print(f"Search Error: {e}")
        # MMR 검색 실패 시 일반 유사도 검색으로 폴백 (더 안정적)
        docs = db.similarity_search(search_query, k=k_val)

    # 4. 컨텍스트 구성 (파일명 태그 포함)
    # AI가 어떤 파일에서 나온 내용인지 알 수 있도록 각 청크에 출처 파일명을 명시
    # 이를 통해 전체 모드에서 모든 강의를 빠짐없이 다이어그램에 반영할 수 있음
    context_chunks = []
    for d in docs:
        source = d.metadata.get('source', '')  # 메타데이터에서 원본 파일명 추출
        if source:
            # OS별 경로 구분자(/ 또는 \) 모두 처리하여 파일명만 추출
            filename = source.split('/')[-1].split('\\')[-1]
            context_chunks.append(f"--- [Source File: {filename}] ---\n{d.page_content}")
        else:
            context_chunks.append(d.page_content)  # 파일명 없으면 내용만 추가

    context = "\n\n".join(context_chunks)  # 청크 사이를 빈 줄로 구분

    # 컨텍스트가 비어있으면 빈 다이어그램 반환
    if not context:
        return 'digraph G { "No Data" [shape=box]; }'

    # 5. 컨텍스트 길이 제한 (토큰 초과 방지)
    # GPT-4o 컨텍스트 윈도우는 128K 토큰이지만, 너무 길면 응답 품질이 저하됨
    # 50,000자 ≈ 15,000~20,000 토큰 (한/영 혼용 기준): 안전한 상한선
    safe_context = context[:50000]

    # 캐시된 LLM 클라이언트 사용 (temperature=0.0: 결정적 출력으로 DOT 코드 일관성 유지)
    llm = get_cached_chat_llm(api_key, temperature=0.0)

    # 한국어 폰트 설정 (Graphviz가 운영체제별로 다른 한글 폰트를 사용)
    font_attr = 'fontname="Malgun Gothic, AppleGothic, sans-serif"'

    # 6. 다이어그램 유형별 Graphviz 스타일 설정
    # 엣지 텍스트 금지: DOT에서 엣지에 레이블이 붙으면 지저분해 보이므로 빈 문자열로 강제
    no_edge_text = 'label="", xlabel="",'

    if "Mind" in viz_type:
        # --- 마인드맵 설정 ---
        layout_engine = "dot"   # 계층형 트리 레이아웃 엔진
        rank_dir = "LR"         # Left-to-Right: 좌→우 방향으로 트리 전개

        # 전체 모드일 때는 강의가 많으므로 노드 간격을 좁혀 한 화면에 더 많이 보이도록 조정
        sep_settings = 'nodesep=0.25; ranksep=0.8;' if is_all_mode else 'nodesep=0.3; ranksep=1.0;'
        graph_attr = f'rankdir={rank_dir}; splines=ortho; {sep_settings} compound=true;'
        # splines=ortho: 엣지를 직각 꺾임선으로 표현 (마인드맵에 적합)
        # compound=true: 클러스터(서브그래프) 간 엣지 허용

        if is_all_mode:
            # 전체 모드: 박스 형태 + 연한 파란 배경 (많은 노드가 정돈되어 보이도록)
            node_def = f'node [shape=box, style="filled,rounded", fillcolor="#E3F2FD", penwidth=1.0, fontsize=12, {font_attr}];'
        else:
            # 상세 모드: 노트 형태 + 연한 노란 배경 (메모처럼 보이는 스타일)
            node_def = f'node [shape=note, style="filled,rounded", fillcolor="#FFF9C4", penwidth=1.0, fontsize=12, margin="0.1,0.1", {font_attr}];'

        # arrowhead=vee: 화살표 머리를 V자 모양으로
        edge_def = f'edge [arrowhead=vee, arrowsize=0.5, color="#546E7A", {no_edge_text} {font_attr}];'

        viz_rules = f"""
        2. **Mindmap Rules**:
            - **Root Node**: Label: **"{root_node}"** (Shape: doubleoctagon, Color: #FFCCBC).
            - **NO EDGE TEXT**: Strictly forbidden. Use plain lines only.
            - **Consistency**: Ensure lecture names correspond to the context provided.
        """
    else:
        # --- 스파이더 다이어그램 설정 ---
        layout_engine = "neato"   # 스프링 모델 기반 자유 배치 엔진 (방사형에 적합)
        graph_attr = 'overlap=false; splines=curved; sep="+25,25"; esep="+10,10"; start=regular;'
        # overlap=false: 노드가 겹치지 않도록
        # splines=curved: 엣지를 곡선으로
        # start=regular: 규칙적인 초기 배치로 방사형 모양 유도

        # plaintext: 노드 테두리 없이 텍스트만 표시 (스파이더 다이어그램에 깔끔)
        node_def = 'node [shape=plaintext, fontcolor="#37474F", fontsize=11, ' + font_attr + '];'
        # arrowhead=none: 방향 없는 엣지 (스파이더는 관계 방향 불필요)
        # len=2.5: 엣지 길이 (노드 간격 조절)
        edge_def = f'edge [arrowhead=none, color="#B0BEC5", len=2.5, penwidth=1.0, {no_edge_text} {font_attr}];'

        viz_rules = f"""
        2. **Spider Diagram Rules**:
            - **Root Node**: Center node **"{root_node}"**.
            - **NO EDGE TEXT**: Strictly forbidden.
        """

    # 7. 최종 프롬프트 구성
    # DOT 코드 템플릿을 제공하여 AI가 올바른 형식으로 출력하도록 유도
    prompt = f"""
    Role: Expert Curriculum Designer & Data Visualization Specialist.
    Task: Generate Graphviz DOT code based on the [Context].
    
    *** VISUALIZATION INSTRUCTION ***
    {scope_instruction}
    
    [Context]
    {safe_context} 
    
    *** STRICT RULES ***
    1. Use ONLY information from the Context.
    2. **Language**: Use the same language as the Context.
    3. **CLEAN EDGES**: **NEVER** put text on edges. Just A -> B.
    4. **NO GENERIC NAMES**: Use the real lecture titles from the [Source File] names or text.
       - FORBIDDEN: "keyword1", "nodeA", "Lecture X".
       - REQUIRED: "Structure", "Thermodynamics", "Lec 2: Metals".
    {viz_rules}
    
    Template:
    digraph G {{
        layout={layout_engine};
        {graph_attr}
        {node_def}
        {edge_def}
        
        // Root Node
        root [label="{root_node}", shape=doubleoctagon, style=filled, fillcolor="#FFCCBC", fontsize=14];
                
        // Define Nodes & Edges
        // ... (Generate nodes for ALL chapters found in context)
    }}
    """

    # 8. LLM 호출 및 비용 추적
    try:
        with get_openai_callback() as cb:
            res = llm.invoke(prompt).content  # DOT 코드 생성
        record_cost("시각화 다이어그램 생성", "gpt-4o",
                    input_tokens=cb.prompt_tokens,
                    output_tokens=cb.completion_tokens)
        return clean_dot_code(res)  # 순수 DOT 코드만 추출하여 반환
    except Exception as e:
        # 에러 발생 시 에러 메시지를 담은 최소한의 유효한 DOT 코드 반환
        # → Streamlit이 빈 문자열 대신 이 코드를 렌더링하여 사용자에게 힌트 제공
        return f'digraph G {{ "Error" [label="Error: {str(e)[:40]}...", shape=box, style=filled, fillcolor="#FFCDD2"]; }}'


# ==============================================================================
# [17] 생성 함수: 플래시카드
# ==============================================================================

def gen_flashcards(db, api_key: str, topic: str, ui_text: dict) -> list:
    """
    벡터 DB에서 내용을 검색하여 Q&A 형태의 플래시카드 리스트를 생성합니다.

    플래시카드 형식:
    [
        {"front": "질문/용어", "back": "정답/설명"},
        ...
    ]

    Args:
        db:       FAISS 벡터 DB
        api_key:  OpenAI API 키
        topic:    주제 필터 (비어있으면 전체 범위)
        ui_text:  현재 언어 팩

    Returns:
        플래시카드 딕셔너리 리스트 (에러 시 에러 카드 1개 포함된 리스트)
    """
    lang = ui_text["target_lang"]
    is_all_mode = check_is_all_mode(topic)

    # 모드별 검색 전략 설정
    if is_all_mode:
        # 전체 모드: 다양한 섹션에서 고루 카드를 생성하기 위해 많은 청크 검색
        search_query = "Important definitions, core concepts, exam questions, summary"
        k_val = 80  # 80개 청크로 전체 강의를 최대한 커버
        scope_instruction = """
        - **Quantity**: **DO NOT LIMIT** the number of cards. Generate as many flashcards as possible to cover the entire context exhaustively.
        - **Scope**: Cover the **ENTIRE breadth** of the provided material from start to finish.
        - **Diversity**: Extract key questions from ALL sections (intro, body, conclusion).
        """
    else:
        # 특정 토픽 모드: 해당 주제 관련 청크만 집중 검색
        search_query = topic
        k_val = 15
        scope_instruction = f"""
        - **Quantity**: Create a comprehensive set of flashcards (no fixed limit) to fully master '{topic}'.
        - **Scope**: Focus **STRICTLY** on the concept of '{topic}'.
        - **Depth**: Ask about definitions, sub-concepts, differences, and applications related specifically to '{topic}'.
        """

    # 벡터 검색 및 컨텍스트 준비
    docs = db.similarity_search(search_query, k=k_val)

    # Fix #1: f-string 내부에서 슬라이싱과 주석을 함께 쓰면 주석이 프롬프트에 포함되는 버그
    # 슬라이싱을 f-string 밖에서 변수로 처리하여 해결
    context = "\n".join([d.page_content for d in docs])
    context_trimmed = context[:15000]  # 컨텍스트 길이 제한 (OpenAI 토큰 한도 고려)

    # 캐시된 OpenAI 클라이언트 사용 (LangChain 없이 직접 호출 → JSON 형식 제어 용이)
    client = get_cached_openai_client(api_key)

    # 프롬프트: JSON 배열만 출력하도록 엄격히 지시
    prompt = f"""
    Role: Exam Prep Tutor.
    Task: Create a comprehensive list of Q&A flashcards based on the [Context].
    
    *** SCOPE INSTRUCTION ***
    {scope_instruction}

    Language: {lang}.
    Format: JSON Array ONLY. Keys: "front" (Question), "back" (Short Answer).
    
    [Context]:
    {context_trimmed}
    
    Output example: [{{"front": "What is X?", "back": "X is Y."}}, {{"front": "...", "back": "..."}}]
    """

    try:
        # system: "You are a JSON generator" → JSON만 출력하도록 역할 부여
        # temperature=0.5: 약간의 다양성으로 질문이 단조롭지 않게
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a JSON generator."},
                {"role": "user",   "content": prompt}
            ],
            temperature=0.5
        )
        # 비용 기록: response.usage로 실제 사용된 토큰 수 확인 (추정이 아닌 정확한 값)
        record_cost("플래시카드 생성", "gpt-4o",
                    input_tokens=response.usage.prompt_tokens,
                    output_tokens=response.usage.completion_tokens)
        res = response.choices[0].message.content
        # clean_json()으로 마크다운 코드블록 등 제거 후 JSON 파싱
        return json.loads(clean_json(res))
    except Exception as e:
        # 파싱 실패 시 에러 카드 1개를 반환하여 UI에서 에러 상황을 표시
        return [{"front": "Error", "back": f"Failed to generate: {str(e)}"}]


# ==============================================================================
# [18] 생성 함수: 퀴즈
# ==============================================================================

def gen_quiz(db, api_key: str, topic: str, ui_text: dict) -> list:
    """
    벡터 DB에서 내용을 검색하여 4지선다 객관식 퀴즈를 생성합니다.

    퀴즈 형식:
    [
        {
            "question":    "질문 텍스트",
            "options":     ["선택지A", "선택지B", "선택지C", "선택지D"],
            "answer":      "선택지A",   # options 중 하나와 정확히 일치하는 문자열
            "explanation": "해설 텍스트"
        },
        ...
    ]

    중요: answer 필드는 options 리스트의 정확한 문자열과 일치해야
    UI에서 ans == q['answer'] 비교가 정확하게 동작합니다.

    Args:
        db:       FAISS 벡터 DB
        api_key:  OpenAI API 키
        topic:    주제 필터
        ui_text:  현재 언어 팩

    Returns:
        퀴즈 딕셔너리 리스트 (에러 시 에러 퀴즈 1개 포함)
    """
    lang = ui_text["target_lang"]
    is_all_mode = check_is_all_mode(topic)  # 전체/특정 토픽 모드 판별

    if is_all_mode:
        # 전체 모드: 다양한 강의/섹션에서 골고루 문제 출제
        search_query = "Exam questions, practice problems, core concepts, critical knowledge"
        k_val = 80
        scope_instruction = """
        - **Quantity**: **DO NOT LIMIT** the number of questions. Generate as many unique questions as possible (e.g., 10, 20, or more) to cover the entire context exhaustively.
        - **Scope**: Questions must cover **various lectures/sections** of the provided material, not just one.
        - **Diversity**: Ensure questions range from fundamental definitions to complex applications found across the entire text.
        """
    else:
        # 특정 토픽 모드: 해당 주제 심층 문제 출제
        search_query = topic
        k_val = 15
        scope_instruction = f"""
        - **Quantity**: Create a comprehensive set of questions (no fixed limit) to fully master '{topic}'.
        - **Scope**: Focus **STRICTLY** on the concept of '{topic}'.
        - **Depth**: Create questions that test the definition, usage, nuances, and common misconceptions of '{topic}' specifically.
        """

    # 벡터 검색 및 컨텍스트 준비
    docs = db.similarity_search(search_query, k=k_val)
    # Fix #1: 주석이 f-string 안으로 들어가는 버그 방지 (플래시카드와 동일한 이유)
    context = "\n".join([d.page_content for d in docs])
    context_trimmed = context[:15000]  # 컨텍스트 길이 제한: 15,000자

    client = get_cached_openai_client(api_key)  # 캐시된 OpenAI 클라이언트 재사용

    prompt = f"""
    Role: Professor.
    Task: Create a comprehensive set of multiple-choice questions based on [Context].
    
    *** SCOPE INSTRUCTION ***
    {scope_instruction}

    Language: {lang}.
    Format: JSON Array ONLY.
    
    Requirements:
    - 4 Options per question.
    - Include clear "explanation" for the correct answer.
    - **Randomize the position of the correct answer** (do not always make 'A' the answer).
    - **IMPORTANT**: The 'answer' field must be the **EXACT String value** from the 'options' list, NOT just 'A', 'B', 'C', or 'D'.
    
    [Context]:
    {context_trimmed}
    
    Output example: 
    [{{"question":"What is 1+1?", "options":["3","2","5","4"], "answer":"2", "explanation":"1+1 equals 2."}}]
    """

    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a JSON generator."},
                {"role": "user",   "content": prompt}
            ],
            temperature=0.5  # 문제 다양성을 위해 약간의 창의성 허용
        )
        record_cost("퀴즈 생성", "gpt-4o",
                    input_tokens=response.usage.prompt_tokens,
                    output_tokens=response.usage.completion_tokens)
        res = response.choices[0].message.content
        return json.loads(clean_json(res))  # 순수 JSON 파싱
    except Exception as e:
        # 에러 퀴즈 반환: UI에서 "question" 키 존재 여부로 성공/실패 판별하므로 키 구조 유지
        return [{"question": "Error", "options": ["Error"], "answer": "Error", "explanation": str(e)}]


# ==============================================================================
# [19] UI 메인 애플리케이션
# ==============================================================================

# --- 세션 상태 초기화 ---
# Streamlit의 session_state: 페이지 새로고침이나 버튼 클릭으로 스크립트가 재실행되어도
# 값이 유지되는 서버 측 변수 저장소 (사용자별로 독립적으로 관리됨)
keys = [
    "chain",          # RAG QA 체인 객체 (get_rag_chain()의 반환값)
    "summary",        # 생성된 요약 마크다운 문자열
    "diagram",        # 생성된 Graphviz DOT 코드 문자열
    "quiz_data",      # 생성된 퀴즈 딕셔너리 리스트
    "flashcards",     # 생성된 플래시카드 딕셔너리 리스트
    "messages",       # AI 도우미 대화 기록 [{"role": "user"/"assistant", "content": "..."}]
    "db",             # FAISS 벡터 DB 객체
    "api_key",        # 입력된 OpenAI API 키 (사이드바에서 입력받아 저장)
    "cost_log",       # API 호출 비용 기록 리스트
    "total_cost_usd", # 세션 누적 비용 (USD)
]
for k in keys:
    if k not in st.session_state:
        st.session_state[k] = None  # 처음 실행 시 모든 키를 None으로 초기화

# None으로 초기화된 것들 중 리스트가 필요한 항목을 빈 리스트/0으로 재설정
if st.session_state.messages is None:       st.session_state.messages = []
if st.session_state.cost_log is None:       st.session_state.cost_log = []
if st.session_state.total_cost_usd is None: st.session_state.total_cost_usd = 0.0


# ==============================================================================
# [20] 사이드바 UI
# ==============================================================================
with st.sidebar:
    # 언어 선택 라디오 버튼 (horizontal=True: 가로 배치)
    lang_opt = st.radio("언어 모드 / Language Mode", ["Korean", "English"], horizontal=True)
    ui = UI[lang_opt]  # 선택된 언어에 해당하는 UI 문자열 딕셔너리 로드

    st.title(ui["sidebar_title"])  # 사이드바 제목 표시

    # API 키 입력 (type="password": 입력 내용이 *로 마스킹됨)
    api_key_input = st.text_input(ui["apikey"], type="password")

    # 파일 업로더 (accept_multiple_files=True: 한 번에 여러 파일 선택 가능)
    # key="lec", key="prob": 두 업로더를 구분하는 내부 식별자
    lec_files  = st.file_uploader(ui["file_label_lec"],  accept_multiple_files=True, key="lec")
    prob_files = st.file_uploader(ui["file_label_prob"], accept_multiple_files=True, key="prob")

    # 분석 시작 버튼
    if st.button(ui["btn_start"], type="primary"):
        # API 키와 파일 중 하나 이상이 있어야 실행
        if api_key_input and (lec_files or prob_files):
            st.session_state.api_key = api_key_input  # API 키를 세션에 저장

            # 지식베이스 구축 (캐싱 적용: 동일 파일이면 API 재호출 없음)
            db = build_knowledge_base(lec_files, prob_files, api_key_input, ui)
            if db:
                st.session_state.db    = db   # FAISS DB 저장
                # RAG 체인 생성 (AI 도우미 탭에서 사용)
                st.session_state.chain = get_rag_chain(db, api_key_input, ui["target_lang"])
                # 파일이 바뀌면 이전 요약/다이어그램 결과를 초기화
                st.session_state.summary = None
                st.session_state.diagram = None
                st.rerun()  # UI 전체 새로고침 (탭 등이 활성화되도록)

    # 비용 패널을 사이드바 하단에 렌더링 (분석 버튼 이후에 표시)
    render_cost_panel()


# ==============================================================================
# [21] 메인 영역 헤더
# ==============================================================================
st.title(ui["title"])                          # 앱 제목
st.markdown(f"**{ui['credit']}**")             # 제작자 정보 (굵게)
st.caption(ui["caption"])                      # 아키텍처 설명 (작은 글씨)


# ==============================================================================
# [22] 탭별 기능 UI (chain과 db가 모두 준비된 경우에만 표시)
# ==============================================================================
if st.session_state.chain and st.session_state.db:
    # 6개 탭 생성: t1~t6에 각각 탭 객체 할당
    t1, t2, t3, t4, t5, t6 = st.tabs(ui["tabs"])


    # ─── 탭 1: 핵심 정리 ────────────────────────────────────────────────────
    with t1:
        # 주제 필터 입력 (label_visibility="collapsed": 레이블 숨김, placeholder만 표시)
        topic_s = st.text_input("Topic_Sum", placeholder=ui["ph_topic"], label_visibility="collapsed")
        if st.button(ui["btn_gen"], key="sum"):
            with st.spinner(ui["spin_gen"]):  # 생성 중 로딩 스피너 표시
                st.session_state.summary = gen_summary(
                    st.session_state.db,
                    st.session_state.api_key,
                    topic_s,
                    ui
                )
        # 생성된 요약이 있으면 마크다운으로 렌더링 (표, 헤더, 볼드 등 지원)
        if st.session_state.summary:
            st.markdown(st.session_state.summary)


    # ─── 탭 2: 시각화 ────────────────────────────────────────────────────────
    with t2:
        # 컨트롤 영역: 유형 선택(2/6) + 주제 입력(3/6) + 생성 버튼(1/6)
        c1, c2, c3 = st.columns([2, 3, 1])
        with c1:
            # 드롭다운 선택: "Mindmap" 또는 "Spider Diagram"
            v_type = st.selectbox("Style", ui["viz_types"], label_visibility="collapsed")
        with c2:
            topic_v = st.text_input("Viz_Topic", placeholder=ui["ph_topic"], label_visibility="collapsed")
        with c3:
            if st.button(ui["btn_gen"], key="viz", use_container_width=True):
                with st.spinner(ui["spin_viz"]):
                    st.session_state.diagram = gen_diagram_optimized(
                        st.session_state.db,
                        st.session_state.api_key,
                        topic_v,
                        v_type,
                        ui
                    )

        # 생성된 DOT 코드가 있으면 렌더링
        if st.session_state.diagram:
            try:
                # use_container_width=True: 컨테이너 너비에 맞게 SVG 크기 조절
                st.graphviz_chart(st.session_state.diagram, use_container_width=True)
                # 디버깅용: DOT 코드 원문을 접기/펼치기 패널에 표시
                with st.expander(ui["err_viz_debug"]):
                    st.code(st.session_state.diagram, language="dot")
            except Exception as e:
                # Graphviz 렌더링 실패 시 에러 메시지 + DOT 코드 원문 표시
                st.error(f"{ui['err_viz']} ({str(e)})")
                st.code(st.session_state.diagram)


    # ─── 탭 3: 플래시카드 ─────────────────────────────────────────────────────
    with t3:
        topic_f = st.text_input("Topic_Flash", placeholder=ui["ph_topic"], label_visibility="collapsed")
        if st.button(ui["btn_gen"], key="flash"):
            with st.spinner(ui["spin_gen"]):
                st.session_state.flashcards = gen_flashcards(
                    st.session_state.db, st.session_state.api_key, topic_f, ui
                )

        # 플래시카드 렌더링
        if st.session_state.flashcards:
            # 정상 데이터 확인: 리스트이고, 첫 항목에 "front" 키가 있는지 검사
            if (isinstance(st.session_state.flashcards, list)
                    and len(st.session_state.flashcards) > 0
                    and "front" in st.session_state.flashcards[0]):
                # 2열 그리드 레이아웃으로 카드 표시
                cols = st.columns(2)
                for i, c in enumerate(st.session_state.flashcards):
                    with cols[i % 2]:  # 짝수 카드는 왼쪽, 홀수 카드는 오른쪽 열
                        st.info(f"**Q{i+1}: {c['front']}**")  # 질문 (파란 배경 박스)
                        with st.expander(ui['lbl_card_back']):  # "정답" 버튼 클릭 시 펼침
                            st.write(c['back'])
            else:
                # 생성 실패 (JSON 파싱 오류 등): 에러 메시지 + 원본 데이터 표시
                st.error(ui["err_json"])
                st.write(st.session_state.flashcards)


    # ─── 탭 4: 퀴즈 ──────────────────────────────────────────────────────────
    with t4:
        topic_q = st.text_input("Topic_Quiz", placeholder=ui["ph_topic"], label_visibility="collapsed")
        if st.button(ui["btn_gen"], key="quiz"):
            with st.spinner(ui["spin_gen"]):
                st.session_state.quiz_data = gen_quiz(
                    st.session_state.db, st.session_state.api_key, topic_q, ui
                )
                # st.rerun(): 라디오 버튼(선택지)의 상태가 올바르게 초기화되도록 즉시 재실행
                st.rerun()

        # 퀴즈 렌더링
        if st.session_state.quiz_data:
            if (isinstance(st.session_state.quiz_data, list)
                    and len(st.session_state.quiz_data) > 0
                    and "question" in st.session_state.quiz_data[0]):
                for i, q in enumerate(st.session_state.quiz_data):
                    st.markdown(f"#### Q{i+1}. {q['question']}")

                    # 라디오 버튼: 선택지 표시 (index=None: 기본값 선택 없음)
                    ans = st.radio(
                        "Select:",
                        q['options'],
                        key=f"q_{i}",       # 각 문제마다 고유한 key 필수 (key 충돌 방지)
                        index=None,          # 처음엔 아무것도 선택되지 않은 상태
                        label_visibility="collapsed"
                    )

                    # 정답 확인 버튼
                    if st.button(ui["quiz_check"], key=f"chk_{i}"):
                        if ans == q['answer']:
                            st.success(ui["quiz_correct"])  # 초록색 성공 메시지
                        else:
                            st.error(ui["quiz_wrong"])      # 빨간색 오류 메시지

                        # 해설 보기 (접기/펼치기)
                        with st.expander(ui["quiz_exp"]):
                            st.write(q['explanation'])
                    st.divider()  # 문제 사이 구분선
            else:
                st.error(ui["err_json"])
                st.write(st.session_state.quiz_data)


    # ─── 탭 5: 오디오 브리핑 ─────────────────────────────────────────────────
    with t5:
        if st.button(ui["btn_gen"], key="audio"):
            # 오디오 브리핑은 요약 텍스트를 기반으로 생성 → 요약이 먼저 있어야 함
            if st.session_state.summary:
                # 캐시된 OpenAI 클라이언트 사용
                client = get_cached_openai_client(st.session_state.api_key)

                try:
                    with st.spinner(ui["spin_audio"]):
                        # 1단계: 요약 텍스트에서 "핵심 내용 요약" 섹션만 추출 + 마크다운 기호 제거
                        core_summary = extract_all_core_parts(st.session_state.summary, ui)

                        # 2단계: OpenAI TTS API의 최대 입력 길이는 4,096자
                        # 안전 마진을 두어 4,000자로 제한 + 잘렸음을 알리는 안내 추가
                        if len(core_summary) > 4000:
                            final_input = core_summary[:4000] + "... (Content truncated due to length limit)"
                            st.caption("⚠️ 텍스트가 너무 길어 앞부분 4000자만 재생됩니다.")
                        else:
                            final_input = core_summary

                        # 3단계: TTS API 호출 (tts-1: 표준 품질, voice="alloy": 중성적 음성)
                        audio = client.audio.speech.create(
                            model="tts-1",
                            voice="alloy",    # 사용 가능한 음성: alloy, echo, fable, onyx, nova, shimmer
                            input=final_input
                        )
                        # 비용 기록: 실제 입력된 글자 수 기준으로 TTS 비용 계산
                        record_cost("오디오 브리핑 (TTS)", "tts-1", chars=len(final_input))

                        st.success("Audio generated! (Reading all 'Key Highlights')")

                        # 4단계: 오디오 플레이어 표시 + 대본 확인 패널
                        st.audio(audio.content, format="audio/mp3")  # 브라우저 내 재생
                        with st.expander("📜 읽어준 대본 (Script)"):
                            st.write(final_input)

                except Exception as e:
                    st.error(f"Error: {str(e)}")
            else:
                # 요약이 없으면 경고 메시지 표시
                st.warning(ui["audio_warn"])


    # ─── 탭 6: AI 도우미 ─────────────────────────────────────────────────────
    with t6:
        # 채팅 히스토리 출력 영역 (height=500: 고정 높이 스크롤 컨테이너)
        chat_box = st.container(height=500)
        for m in st.session_state.messages:
            chat_box.chat_message(m["role"]).write(m["content"])

        # st.chat_input(): 화면 하단에 고정된 채팅 입력창
        # :=  (왈러스 연산자): 입력값을 q에 할당하면서 동시에 참/거짓 판별
        if q := st.chat_input(ui["chat_ph"]):
            # 사용자 메시지를 히스토리에 추가 + 즉시 화면에 표시
            st.session_state.messages.append({"role": "user", "content": q})
            chat_box.chat_message("user").write(q)

            with st.spinner("Analyzing intent & Searching documents..."):
                try:
                    # === 질문 의도 분석 및 검색 쿼리 확장 ===
                    # 사용자의 짧은 질문을 LLM이 더 잘 이해할 수 있도록 의도에 맞게 확장
                    q_lower = q.strip().lower()

                    if q_lower in ["all", "전부", "전체", "everything", "요약해줘"]:
                        # [전체 요약 모드]: 문서 전체를 아우르는 포괄적 요약 요청
                        search_query = (
                            "Provide a comprehensive and very detailed summary of the ENTIRE provided material. "
                            "Cover all lectures, core concepts, structure, and main arguments from start to finish. "
                            "Do not miss any major sections."
                        )

                    elif any(x in q_lower for x in ["시험", "exam", "test", "유형", "type", "strategy", "plan", "계획", "대비"]):
                        # [시험 전략 모드]: 연습 문제 태그([Practice Problem])가 붙은 내용을
                        # 분석하여 시험 유형과 학습 전략을 제시
                        search_query = (
                            f"User Question: '{q}'\n\n"
                            "Task: Act as an Exam Strategist. "
                            "1. Analyze the content labeled `[Practice Problem]` to identify exam styles (MCQ, Essay, etc.) and difficulty. "
                            "2. Summarize the types of questions that appear. "
                            "3. Provide a concrete study plan and preparation strategy based on these patterns."
                        )

                    else:
                        # [일반 개념 질문 모드]: 특정 용어나 개념에 대한 상세 설명 요청
                        search_query = (
                            f"Explain the concept of '{q}' in great detail. "
                            "Include its definition, context, related terms, and why it is important in this document."
                        )

                    # RAG 체인 실행: 위에서 구성한 확장된 search_query를 {question} 자리에 삽입
                    # → FAISS 검색 + GPT-4o 답변 생성이 한 번에 처리됨
                    with get_openai_callback() as cb:
                        response = st.session_state.chain.invoke({"query": search_query})
                    record_cost("AI 도우미 답변", "gpt-4o",
                                input_tokens=cb.prompt_tokens,
                                output_tokens=cb.completion_tokens)
                    res = response['result']  # RetrievalQA 결과에서 답변 텍스트 추출

                except Exception as e:
                    res = f"Error: {str(e)}"

            # AI 응답을 채팅 UI에 표시 + 히스토리에 저장
            chat_box.chat_message("assistant").write(res)
            st.session_state.messages.append({"role": "assistant", "content": res})


# DB가 아직 없으면 사이드바에서 파일을 업로드하도록 안내
else:
    st.info(f"👈 {ui['sidebar_title']}")
